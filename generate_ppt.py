import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Configure slide size to widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Colors
    DARK_BG = RGBColor(15, 23, 42)      # Slate-900 (Background)
    LIGHT_GRAY = RGBColor(226, 232, 240) # Slate-200 (Primary Text)
    MUTED_GRAY = RGBColor(148, 163, 184) # Slate-400 (Secondary Text)
    ACCENT_BLUE = RGBColor(56, 189, 248)  # Sky-400 (Input / Injection highlights)
    ACCENT_PURPLE = RGBColor(192, 132, 252) # Purple-400 (Output / Hallucination highlights)
    ACCENT_TEAL = RGBColor(45, 212, 191)   # Teal-400 (Status / Training highlights)
    CARD_BG = RGBColor(30, 41, 59)      # Slate-800 (Card background)
    WHITE = RGBColor(255, 255, 255)
    
    blank_layout = prs.slide_layouts[6] # Blank layout for full custom drawing

    def apply_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    def add_title(slide, text, color=WHITE):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_left = Inches(0)
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color
        return title_box

    def add_bullet_points(slide, points, left, top, width, height, font_size=16):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        
        for idx, pt in enumerate(points):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = pt
            p.font.name = 'Segoe UI'
            p.font.size = Pt(font_size)
            p.font.color.rgb = LIGHT_GRAY
            p.space_after = Pt(10)
            p.level = 0
        return box

    def draw_card(slide, left, top, width, height, title="", text="", line_color=ACCENT_BLUE, fill_color=CARD_BG):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
        
        if title or text:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.2)
            tf.margin_right = Inches(0.2)
            tf.margin_top = Inches(0.2)
            
            p = tf.paragraphs[0]
            if title:
                p.text = title
                p.font.name = 'Segoe UI'
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.alignment = PP_ALIGN.LEFT
                
                if text:
                    p2 = tf.add_paragraph()
                    p2.text = text
                    p2.font.name = 'Segoe UI'
                    p2.font.size = Pt(14)
                    p2.font.color.rgb = LIGHT_GRAY
                    p2.space_before = Pt(8)
                    p2.alignment = PP_ALIGN.LEFT
            else:
                p.text = text
                p.font.name = 'Segoe UI'
                p.font.size = Pt(15)
                p.font.color.rgb = LIGHT_GRAY
                p.alignment = PP_ALIGN.CENTER
        return shape

    def draw_arrow(slide, left, top, width, height, color=MUTED_GRAY, direction='right'):
        shape_type = MSO_SHAPE.RIGHT_ARROW if direction == 'right' else MSO_SHAPE.DOWN_ARROW
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_background(slide1)
    
    # Title & Subtitle centered
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "PROMPT INJECTION & HALLUCINATION\nGUARDRAIL SYSTEM"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(20)
    
    p2 = tf.add_paragraph()
    p2.text = "A Dual-Layer Defense Shield for LLM Applications"
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(22)
    p2.font.color.rgb = ACCENT_BLUE
    p2.alignment = PP_ALIGN.CENTER
    
    # Footer metadata
    meta_box = slide1.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(11.33), Inches(0.8))
    tf_meta = meta_box.text_frame
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = "Backend API Security  |  Output Quality Assurance  |  LangChain Sandbox"
    p_meta.font.name = 'Segoe UI'
    p_meta.font.size = Pt(13)
    p_meta.font.color.rgb = MUTED_GRAY
    p_meta.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 2: The Core Problem
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_background(slide2)
    add_title(slide2, "The Core Problem: Vulnerabilities in LLM Apps")
    
    # Left Box: Prompt Injection
    draw_card(slide2, Inches(1.0), Inches(1.8), Inches(5.0), Inches(4.5), 
              title="1. Input Vulnerability: Prompt Injection", 
              text="Attackers hide adversarial instructions inside normal-looking messages to hijack the LLM.\n\n"
                   "• System Hijack: Overrides system rules ('Ignore previous rules...').\n"
                   "• Data Leak: Pulls system passwords or database context.\n"
                   "• Uncontrolled API Calls: Leads to unauthorized operations.", 
              line_color=RGBColor(251, 146, 60))
    
    # Right Box: Hallucination
    draw_card(slide2, Inches(7.33), Inches(1.8), Inches(5.0), Inches(4.5), 
              title="2. Output Vulnerability: LLM Hallucinations", 
              text="Models generate incorrect, factually wrong, or ungrounded responses unsupported by the source database.\n\n"
                   "• Factual Contradictions: Lies or alters context variables.\n"
                   "• Silent Failures: Plausible-sounding false statements.\n"
                   "• User Distrust: Causes incorrect advice or brand damage.", 
              line_color=ACCENT_PURPLE)

    # ==========================================
    # SLIDE 3: System Architecture
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_background(slide3)
    add_title(slide3, "System Architecture: End-to-End Pipeline")
    
    # Horizontal flow blocks
    block_width = Inches(1.8)
    block_height = Inches(2.2)
    top_pos = Inches(2.2)
    arrow_width = Inches(0.4)
    arrow_height = Inches(0.3)
    arrow_top = Inches(3.15)
    
    # User message
    draw_card(slide3, Inches(0.75), top_pos, block_width, block_height, 
              title="1. User Input", text="Sends prompt query to app.", line_color=WHITE)
    
    draw_arrow(slide3, Inches(2.7), arrow_top, arrow_width, arrow_height)
    
    # Guardrail Input Filter
    draw_card(slide3, Inches(3.25), top_pos, block_width, block_height, 
              title="2. Input Shield", text="Regex rules + DistilBERT model.", line_color=ACCENT_BLUE)
              
    draw_arrow(slide3, Inches(5.2), arrow_top, arrow_width, arrow_height)
    
    # ChatGPT Core
    draw_card(slide3, Inches(5.75), top_pos, block_width, block_height, 
              title="3. Target LLM", text="Processes safe prompts (ChatGPT).", line_color=ACCENT_TEAL)
              
    draw_arrow(slide3, Inches(7.7), arrow_top, arrow_width, arrow_height)
    
    # Guardrail Output Scorer
    draw_card(slide3, Inches(8.25), top_pos, block_width, block_height, 
              title="4. Output Scorer", text="NLI + Embedding similarity checks.", line_color=ACCENT_PURPLE)
              
    draw_arrow(slide3, Inches(10.2), arrow_top, arrow_width, arrow_height)
    
    # Safe Output / Alert block
    draw_card(slide3, Inches(10.75), top_pos, block_width, block_height, 
              title="5. Response", text="Renders safe content or block alert.", line_color=WHITE)
              
    # Summary list at the bottom
    bullets = [
        "✔ Dual-layered security shield wrapping standard LangChain execution pipelines.",
        "✔ Low-latency backend Flask API handles high-speed inspection calls before LLM triggers."
    ]
    add_bullet_points(slide3, bullets, Inches(0.75), Inches(5.0), Inches(11.83), Inches(1.5))

    # ==========================================
    # SLIDE 4: Prompt Injection Guardrail (Input Shield)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_background(slide4)
    add_title(slide4, "Input Protection: Prompt Injection Detector")
    
    # Left column: details
    bullets = [
        "1. Heuristic Regex Rules (Fast Path): Scans for known malicious prompt signatures ('ignore guidelines', 'developer mode', 'dan bypass') immediately.",
        "2. Custom DistilBERT Classifier (AI Path): Processed by a sequence classification model trained to calculate adversarial threat probability.",
        "3. Decisive Action: Any threat score exceeding 0.50 terminates pipeline flow immediately, returning a secure warning instead of invoking ChatGPT."
    ]
    add_bullet_points(slide4, bullets, Inches(0.75), Inches(1.8), Inches(6.0), Inches(4.5))
    
    # Right column: flowchart
    box_w = Inches(3.5)
    box_h = Inches(0.7)
    flow_left = Inches(8.0)
    
    draw_card(slide4, flow_left, Inches(1.8), box_w, box_h, text="User Input Prompt", line_color=WHITE)
    draw_arrow(slide4, flow_left + Inches(1.6), Inches(2.6), Inches(0.3), Inches(0.4), direction='down')
    
    draw_card(slide4, flow_left, Inches(3.1), box_w, box_h, text="Rule Heuristics Check", line_color=ACCENT_BLUE)
    draw_arrow(slide4, flow_left + Inches(1.6), Inches(3.9), Inches(0.3), Inches(0.4), direction='down')
    
    draw_card(slide4, flow_left, Inches(4.4), box_w, box_h, text="DistilBERT Classifier", line_color=ACCENT_BLUE)
    draw_arrow(slide4, flow_left + Inches(1.6), Inches(5.2), Inches(0.3), Inches(0.4), direction='down')
    
    draw_card(slide4, flow_left, Inches(5.7), box_w, box_h, text="Blocked (Alert) OR Allowed (Route)", line_color=ACCENT_TEAL)

    # ==========================================
    # SLIDE 5: Hallucination Guardrail (Output Shield)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_background(slide5)
    add_title(slide5, "Output Protection: Consistency & Alignment Scorer")
    
    # Left column details
    bullets = [
        "Dual-method evaluation ensures robust grounding check:",
        "1. DeBERTa NLI Cross-Encoder: Formulates context as Premise and reply as Hypothesis. Classifies if reply is logical Entailment (grounded) vs. Contradiction (hallucination).",
        "2. Sentence-level Embedding Similarity (SBERT): Splits LLM output into sentences, encoding and comparing each sentence against the grounding context.",
        "3. Ensemble Decision: Averages NLI and Semantic scores. If aggregate grounding falls below the 55% threshold, the output is blocked and marked ungrounded."
    ]
    add_bullet_points(slide5, bullets, Inches(0.75), Inches(1.8), Inches(6.0), Inches(4.8))
    
    # Right column: Diagram
    draw_card(slide5, Inches(7.5), Inches(1.8), Inches(4.5), Inches(1.0), 
              title="Grounding Facts (Context) + LLM Response", line_color=WHITE)
              
    draw_arrow(slide5, Inches(8.3), Inches(2.9), Inches(0.3), Inches(0.3), direction='down')
    draw_arrow(slide5, Inches(10.8), Inches(2.9), Inches(0.3), Inches(0.3), direction='down')
    
    # Parallel pathways
    draw_card(slide5, Inches(7.2), Inches(3.3), Inches(2.2), Inches(1.2), 
              title="NLI Model", text="Entailment / Contradiction", line_color=ACCENT_PURPLE)
              
    draw_card(slide5, Inches(9.8), Inches(3.3), Inches(2.2), Inches(1.2), 
              title="SBERT Model", text="Cosine Similarity", line_color=ACCENT_PURPLE)
              
    draw_arrow(slide5, Inches(8.3), Inches(4.6), Inches(0.3), Inches(0.3), direction='down')
    draw_arrow(slide5, Inches(10.8), Inches(4.6), Inches(0.3), Inches(0.3), direction='down')
    
    # Join Block
    draw_card(slide5, Inches(7.5), Inches(5.0), Inches(4.5), Inches(1.2), 
              title="Ensemble Score Averager", text="Calculates combined alignment score. Decides safe vs. block.", line_color=ACCENT_TEAL)

    # ==========================================
    # SLIDE 6: Dataset Expansion & Perturbation Engine
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_background(slide6)
    add_title(slide6, "Data Source: Factual Perturbation & Expansion")
    
    bullets = [
        "1. Base Factual Reference: Baseline grounding contexts loaded in database.",
        "2. Heuristic Factual Perturbation: Generates negative data (hallucinations) programmatically by altering contexts:",
        "    • Entity Swapping (e.g., 'Paris' ➔ 'Berlin')",
        "    • Verb Negation (e.g., 'boils' ➔ 'does not boil')",
        "    • Numeric Shifts (e.g., '100°C' ➔ '200°C')",
        "3. Generative Adversarial Prompts: Uses LLM to synthetically generate advanced prompt injection styles (virtualization, translation).",
        "4. Model Fine-Tuning: Continually re-trains DistilBERT in a background Python thread to fit the expanded, custom security dataset."
    ]
    add_bullet_points(slide6, bullets, Inches(0.75), Inches(1.8), Inches(6.0), Inches(5.0))
    
    # Cycle flow diagram on right
    cw = Inches(4.0)
    ch = Inches(0.9)
    cx = Inches(7.8)
    
    draw_card(slide6, cx, Inches(1.8), cw, ch, title="1. Baseline Factual Context", line_color=WHITE)
    draw_arrow(slide6, cx + Inches(1.85), Inches(2.8), Inches(0.3), Inches(0.3), direction='down')
    
    draw_card(slide6, cx, Inches(3.2), cw, ch, title="2. Programmatic Perturbation", text="Verb negation / entity swapping", line_color=ACCENT_BLUE)
    draw_arrow(slide6, cx + Inches(1.85), Inches(4.2), Inches(0.3), Inches(0.3), direction='down')
    
    draw_card(slide6, cx, Inches(4.6), cw, ch, title="3. Synthetic Generator (OpenAI)", text="Generates diverse attacks", line_color=ACCENT_BLUE)
    draw_arrow(slide6, cx + Inches(1.85), Inches(5.6), Inches(0.3), Inches(0.3), direction='down')
    
    draw_card(slide6, cx, Inches(6.0), cw, ch, title="4. Retrained Guardrail Model", line_color=ACCENT_TEAL)

    # ==========================================
    # SLIDE 7: LangChain Integration Flow
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_background(slide7)
    add_title(slide7, "LangChain Production Integration")
    
    bullets = [
        "Seamless deployment wrapper fits existing AI applications:",
        "• Unified Runnable Wrapper: The system wraps standard ChatGPT endpoints inside a single transaction wrapper.",
        "• Custom Interception: Captures inputs before hitting ChatGPT, and intercepts outputs before rendering to the client.",
        "• Low Latency: Local neural networks run inference instantly (CPU or GPU).",
        "• Security Logging: All transaction parameters and block alerts are piped to a central real-time Guardrail Auditor."
    ]
    add_bullet_points(slide7, bullets, Inches(0.75), Inches(1.8), Inches(6.0), Inches(4.8))
    
    # Process sequence shapes
    pw = Inches(4.2)
    ph = Inches(0.9)
    px = Inches(7.8)
    
    draw_card(slide7, px, Inches(1.8), pw, ph, title="1. Intercept User Prompt", text="Scans inputs for injections.", line_color=ACCENT_BLUE)
    draw_arrow(slide7, px + Inches(1.95), Inches(2.8), Inches(0.3), Inches(0.3), direction='down')
    
    draw_card(slide7, px, Inches(3.2), pw, ph, title="2. Invoke Target LLM", text="Sends only verified safe prompts to ChatGPT.", line_color=ACCENT_TEAL)
    draw_arrow(slide7, px + Inches(1.95), Inches(4.2), Inches(0.3), Inches(0.3), direction='down')
    
    draw_card(slide7, px, Inches(4.6), pw, ph, title="3. Intercept LLM Response", text="Extracts raw response and compares against facts.", line_color=ACCENT_BLUE)
    draw_arrow(slide7, px + Inches(1.95), Inches(4.2), Inches(0.3), Inches(0.3), direction='down')
    
    draw_card(slide7, px, Inches(6.0), pw, ph, title="4. Factual Output Validation", text="Blocks hallucinations or yields safe text.", line_color=ACCENT_PURPLE)

    # ==========================================
    # SLIDE 8: Verification & Testing Results
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    apply_background(slide8)
    add_title(slide8, "Verification, Testing & Results")
    
    # Visual grid of metric cards
    w_card = Inches(5.4)
    h_card = Inches(2.0)
    
    # Injection results card
    draw_card(slide8, Inches(0.75), Inches(1.8), w_card, h_card,
              title="Prompt Injection Scanner Validation",
              text="✔ 100% of seed adversarial injection attacks correctly identified.\n"
                   "✔ Regex catch-all intercepts common keywords immediately.\n"
                   "✔ DistilBERT classifier accurately scores virtualization and roleplay attacks.",
              line_color=ACCENT_BLUE)
              
    # Hallucination results card
    draw_card(slide8, Inches(7.18), Inches(1.8), w_card, h_card,
              title="Hallucination Scorer Validation",
              text="✔ Cleanly distinguishes perturbed context lies from aligned truths.\n"
                   "✔ Average grounding checks minimize false positive classifications.\n"
                   "✔ Sentences with <55% cosine similarity are successfully blocked.",
              line_color=ACCENT_PURPLE)
              
    # Integration and speed card
    draw_card(slide8, Inches(0.75), Inches(4.3), w_card, h_card,
              title="Integration & Performance",
              text="✔ LangChain transaction wrappers securely intercept LLM calls.\n"
                   "✔ Offline simulation mode runs seamlessly without an API key.\n"
                   "✔ Dashboard auditor prints real-time safety indices and warning details.",
              line_color=ACCENT_TEAL)
              
    # Future expansions card
    draw_card(slide8, Inches(7.18), Inches(4.3), w_card, h_card,
              title="Model Fine-Tuning & Datasets",
              text="✔ Flask background threads manage retraining pipelines in PyTorch.\n"
                   "✔ Expanded datasets update local classifier weights in './models/'.\n"
                   "✔ Dashboard allows click-to-generate synthetic data scaling.",
              line_color=WHITE)

    # Save presentation
    output_filename = "presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation successfully created and saved to: {os.path.abspath(output_filename)}")

if __name__ == "__main__":
    create_presentation()
